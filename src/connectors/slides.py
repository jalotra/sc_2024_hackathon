from .text import Parse
from pathlib import Path
from pptx import Presentation
import json


class PptParser(Parse):
    def parse(self, blobOrFileName) -> str:
        if self.is_filename(blobOrFileName):
            if not Path(blobOrFileName).exists():
                raise FileNotFoundError(
                    f"Not able to find the file at path : {blobOrFileName}"
                )
            prs = Presentation(blobOrFileName)

            # text_runs will be populated with a list of strings,
            # one for each text run in presentation
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)

            return json.dumps(text_runs)

        raise Exception(
            f"Something bad happened; fileName needed; got this : {blobOrFileName}"
        )
